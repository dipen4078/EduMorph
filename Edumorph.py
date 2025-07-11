from flask import Flask, render_template, request, redirect, url_for, session
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user

from groq import Groq
from fpdf import FPDF
import shutil
import re
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import os
import subprocess
from pdf2image import convert_from_path
from pptx.util import Pt  # For setting font size
from pptx.enum.text import PP_ALIGN  # For alignment
import json


app = Flask(__name__)
app.config['SECRET_KEY'] = 'your_secret_key' # no need for api key
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///users.db'
db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

# Initialize Groq with Llama3
groq = Groq(api_key='your_api_key')   # Detail about generating the api key is in the repo file named generate_apikey.txt


class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    password = db.Column(db.String(150), nullable=False)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        password = bcrypt.generate_password_hash(request.form['password']).decode('utf-8')
        user = User(username=username, password=password)
        db.session.add(user)
        db.session.commit()
        return redirect(url_for('login'))
    return render_template('register.html')


@app.route('/login', methods=['GET', 'POST'])
def login():
    error=None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        user = User.query.filter_by(username=username).first()
        if user and bcrypt.check_password_hash(user.password, password):
            login_user(user)
            return redirect(url_for('dashboard'))
        else:
            error = "Incorrect Credentials. Please try again."

    return render_template('login.html',error=error)


@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('index'))


@app.route('/dashboard')
@login_required
def dashboard():
    return render_template('dashboard.html')


@app.route('/select_topic', methods=['GET', 'POST'])
@login_required
def select_topic():
    if request.method == 'POST':
        topic = request.form['topic']
        level = request.form['level']
        return redirect(url_for('list_subtopics', topic=topic, level=level))
    return render_template('select_topic.html')


@app.route('/list_subtopics/<topic>/<level>', methods=['GET'])
@login_required
def list_subtopics(topic, level):
    subtopics_lines = get_subtopics(topic,level)
    # subtopics_lines = subtopics_str.splitlines()
    return render_template('list_subtopics.html', subtopics=subtopics_lines, topic=topic, level=level)


@app.route('/learn_topic/<topic>/<subtopic>/<level>')
@login_required
def learn_topic(topic, subtopic, level):
    resources = get_learning_resources(topic, subtopic)
    roadmap = generate_roadmap(topic, subtopic)
    
    pdf_filename = generate_pdf(subtopic,topic,level)

    return render_template('learn_topic.html', topic=topic, subtopic=subtopic, level=level, resources=resources, roadmap=roadmap, pdf_filename=pdf_filename)


def get_subtopics(topic,level):
    # prompt = f"Generate subtopics for the topic: {topic} for a {level}"
    prompt = f"""
    Generate subtopics for the topic: {topic} for a {level}.
    Provide the output in the following format:
    1. Main Topic 1
        - Subtopic 
        - Subtopic 
    2. Main Topic 2
        - Subtopic 
        - Subtopic 
    """
    chat_completion = groq.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama3-8b-8192",
    )
    response = chat_completion.choices[0].message.content
    # Process the response to remove '*' and properly format the text
    lines = response.splitlines()
    formatted_subtopics = []
    
    for line in lines:
        # Remove any '*' from the line and clean it
        clean_line = line.replace('*', '').strip()
        if clean_line:
            formatted_subtopics.append(clean_line)

    return formatted_subtopics

def get_learning_resources(topic, subtopic):
    prompt = f"Provide a learning roadmap for the subtopic '{subtopic}' under the topic '{topic}'. Include resources like articles, tutorials, and videos."
    chat_completion = groq.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama3-8b-8192",
    )
    response = chat_completion.choices[0].message.content
    # Remove asterisks from the response
    cleaned_response = response.replace('*', '')

    return cleaned_response

def generate_roadmap(topic, subtopic):
    # Placeholder logic to generate a roadmap based on topic and subtopic
    roadmap = f"To learn {subtopic} under {topic}, consider the following steps:\n"
    roadmap += "1. Start with the basics.\n"
    roadmap += "2. Practice through examples.\n"
    roadmap += "3. Explore advanced topics.\n"
    roadmap += "4. Use resources like: \n"
    roadmap += "   - Online tutorials\n"
    roadmap += "   - Videos\n"
    roadmap += "   - Documentation\n"
    return roadmap


def create_video(slides_folder, audio_folder, output_video, fps=24):
    slides = sorted([file for file in os.listdir(slides_folder) if file.endswith(('.png', '.jpg', '.jpeg'))])
    audio_files = sorted([file for file in os.listdir(audio_folder) if file.endswith('.mp3')])

    if len(slides) != len(audio_files):
        raise Exception("Number of slides and audio files do not match.")

    clips = []
    for slide, audio in zip(slides, audio_files):
        slide_path = os.path.join(slides_folder, slide)
        audio_path = os.path.join(audio_folder, audio)

        # Create an ImageClip for the slide
        image_clip = ImageClip(slide_path).set_duration(5)  # Default duration

        # Load the audio clip
        audio_clip = AudioFileClip(audio_path)

        # Set the duration of the image clip to match the audio
        image_clip = image_clip.set_duration(audio_clip.duration)

        # Set the audio of the image clip
        image_clip = image_clip.set_audio(audio_clip)

        clips.append(image_clip)

    # Concatenate all clips into one video
    final_video = concatenate_videoclips(clips, method="compose")

    # Write the final video to a file
    final_video.write_videofile(output_video, fps=fps)

    print(f"Video saved to {output_video}")


def convert_pptx_to_images(pptx_path, images_folder):
    # Ensure output folder exists
    os.makedirs(images_folder, exist_ok=True)

    # Convert PPTX to PDF using LibreOffice
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    try:
        subprocess.run([
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "--headless", "--convert-to", "pdf", pptx_path, "--outdir", os.path.dirname(pptx_path)
        ], check=True)

        # Convert PDF pages to images
        images = convert_from_path(pdf_path)
        for i, image in enumerate(images, start=1):
            image_path = os.path.join(images_folder, f"slide_{i}.png")
            image.save(image_path, "PNG")
            print(f"Exported Slide {i} to {image_path}")

    except Exception as e:
        print("Error during conversion:", e)

    finally:
        # Clean up intermediate PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)


def generate_voice_explanations(slide_contents, audio_folder, language='en'):
    os.makedirs(audio_folder, exist_ok=True)
    for idx, slide in enumerate(slide_contents):
        prompt = f" Give explanation for the following briefly: Slide {idx + 1}: {slide.get('title', '')}. " + " ".join(slide.get('bullet_points', []))
        chat_completion = groq.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama3-8b-8192",
        )
        explanation_text = chat_completion.choices[0].message.content
        explanation_text=explanation_text.replace('*',"")
        explanation_text=explanation_text.replace('"',"")
        explanation_text=explanation_text.replace("/","")
        explanation_text=explanation_text.replace("\\","")
        explanation_text=explanation_text.replace("'","")
        explanation_text=explanation_text.replace("]","")
        explanation_text=explanation_text.replace("[","")
        explanation_text=explanation_text.replace("(","")
        explanation_text=explanation_text.replace(")","")
        explanation_text=explanation_text.replace("}","")
        explanation_text=explanation_text.replace("{","")
        explanation_text=explanation_text.replace("!","")
        explanation_text=explanation_text.replace("$","")
        explanation_text=explanation_text.replace("`","")
        explanation_text=explanation_text.replace(":","")
        explanation_text=explanation_text.replace(";","")
        explanation_text=explanation_text.replace("~","")
        #explanation_text = f"Slide {idx + 1}: {slide.get('title', '')}. " + " ".join(slide.get('bullet_points', []))
        tts = gTTS(text=explanation_text, lang=language)
        audio_path = os.path.join(audio_folder, f"slide_{idx + 1}.mp3")
        tts.save(audio_path)
        print(f"Audio saved to {audio_path}")

def summarize_content(content):
    """
    Parses the input content string and splits it into slides based on slide headings.

    Args:
        content (str): The input text containing slide information.

    Returns:
        list: A list of slide dictionaries with 'title' and 'bullet_points'.
    """
    # Regular expression to match "Slide X:" followed by any content (without capturing groups)
    slide_pattern = re.compile(r'Slide\s+\d+:\s*.*?(?=Slide\s+\d+:|$)', re.DOTALL | re.IGNORECASE)

    # Find all slide blocks
    matches = slide_pattern.findall(content)

    slides = []
    for match in matches:
        # Split the match into lines
        lines = match.strip().split('\n')
        if not lines:
            continue

        # Extract the title from the first line
        title_line = lines[0].strip()
        title_match = re.match(r'Slide\s+\d+:\s*(.*)', title_line, re.IGNORECASE)
        if title_match:
            title = title_match.group(1).strip()
        else:
            title = "Slide"  # Default title if none is found

        # Extract bullet points from the remaining lines
        bullet_points = []
        for line in lines[1:]:
            # Remove leading symbols like '*', '-', '•', or any whitespace
            bullet = re.sub(r'^[\-\*\•\s]*', '', line).strip()
            if bullet:
                bullet_points.append(bullet)

        slides.append({
            'title': title,
            'bullet_points': bullet_points
        })

    return slides


def print_slides(slides):
    for slide in slides:
        print(slide['title'])  # Print the slide title
        for bullet in slide['bullet_points']:
            print(f" - {bullet}")  # Print each bullet point
        print()  # Add a newline between slides for better readability

def create_ppt(slides, output_file):
    prs = Presentation()

    for slide in slides:
        slide_layout = prs.slide_layouts[1]  # Use the slide layout with title and content
        slide_obj = prs.slides.add_slide(slide_layout)

        # Use the first bullet point as the title of the slide
        title_text = slide['title'] 
        i = 0
        while(title_text == "" and i<len(slide['bullet_points'])):
            title_text = slide['bullet_points'].pop(i) if slide['bullet_points'] else "Untitled Slide"
        if(title_text == ""):
            title_text = "Untitled Slide"
        # Set the title

        title = slide_obj.shapes.title
        title.text = title_text

        # Format the title
        title_frame = title.text_frame
        p = title_frame.paragraphs[0]  # The title text is the first paragraph
        run = p.runs[0]

        run.font.bold = True  # Bold
        run.font.italic = True  # Italic
        run.font.underline = True  # Underline
        run.font.name = 'Adobe Devanagari'  # Set font
        run.font.size = Pt(24)  # Set font size to 24 pt

        # Set title alignment to center
        p.alignment = PP_ALIGN.CENTER

        # Add remaining bullet points (left aligned, font size 18, no special formatting)
        content = slide_obj.shapes.placeholders[1].text_frame
        content.clear()  # Clear the default content placeholder

        # Leave two blank lines
        blank_para = content.add_paragraph()
        # blank_para.text = "\n"

        for bullet in slide['bullet_points']:
            p = content.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)  # Set font size for bullet points
            p.font.bold = False
            p.font.underline = False
            p.font.italic = False
            p.font.name = 'Arial'  # Set a different font if desired
            p.alignment = PP_ALIGN.LEFT  # Left aligned bullet points

    prs.save(output_file)
    print(f"PowerPoint saved to {output_file}")



@app.route('/generate_slides/<topic>/<subtopic>/<level>', methods=['GET'])
@login_required
def generate_slides(topic, subtopic, level):
    slides = create_slides(topic, subtopic)
    
    send_string = ""
    
    for i in range(0, len(slides)):
        send_string += slides[i] + '\n'

    sli = summarize_content(send_string)
    # print_slides(sli)

    # Create PowerPoint from slides
    output_ppt_file = os.path.abspath(f"presentation.pptx")
    create_ppt(sli, output_ppt_file)

    # Verify the PowerPoint file was created
    if not os.path.exists(output_ppt_file):
        raise FileNotFoundError(f"The file {output_ppt_file} does not exist.")

    # Generate audio explanations
    audios_folder = os.path.abspath("audio_explanations")
    generate_voice_explanations(sli, audios_folder)

    # Convert PPTX to images
    images_folder = os.path.abspath("slides_images")
    convert_pptx_to_images(output_ppt_file, images_folder)
    
    # Create the video
    create_video(images_folder, audios_folder, os.path.join('static', 'presentation_video.mp4'))
    # create_video(images_folder, audios_folder, "presentation_video.mp4")

    # Delete the audio and images folders after creating the video
    if os.path.exists(audios_folder):
        shutil.rmtree(audios_folder)
    if os.path.exists(images_folder):
        shutil.rmtree(images_folder)
    
    # Redirect to video playback
    return redirect(url_for('play_video', topic=topic, subtopic=subtopic))

@app.route('/play_video/<topic>/<subtopic>')
@login_required
def play_video(topic, subtopic):
    return render_template('play_video.html', topic=topic, subtopic=subtopic)


def create_slides(topic, subtopic):
    # Use LangChain to create slides. Here is a placeholder for the logic.
    prompt = f'''Create a presentation on {subtopic} under the topic {topic}. Make sure that the slide contains detailed information regarding the topic and also includes necessary coding snippet titled "Code:" if the topic is coding related.
    The slide format is as follows -> The first line should contain the "slide number Topic, From the next line it should give Several bullets explaining the topic which should contain only text without any code examples. '''
        
    # Assuming you have set up LangChain properly:
    slides_generation = groq.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama3-8b-8192",
    )
    
    slides = slides_generation.choices[0].message.content.splitlines()
    
    # Remove all asterisks from each slide
    cleaned_slides = [slide.replace('*', '') for slide in slides]
    
    return cleaned_slides


def generate_pdf(subtopic, topic, level):
    # Query Groq for detailed learning objectives and explanations
    prompt = f"Explain the concept of {subtopic} in a detailed manner for a {level}. Start by introducing the key concepts, followed by step-by-step explanations. Provide at least two examples (one simple, one slightly more complex). For programming concepts, include code snippets and explain what each part of the code does."
    chat_completion = groq.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama3-8b-8192",
    )
    response = chat_completion.choices[0].message.content
    cleaned_response = response.replace('*', '')

    # Initialize PDF
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Set Title
    pdf.set_font('Arial', 'BI', 16)  # Bold and Italics
    title_text = f'Learning Roadmap for {subtopic} under {topic}'
    
    # Add title
    pdf.cell(200, 10, title_text, ln=True, align='C')

    # Underline the title manually by drawing a line under it
    title_width = pdf.get_string_width(title_text) + 6  # Get the width of the title
    pdf.set_x((210 - title_width) / 2)  # Center the line under the text
    pdf.cell(title_width, 0, '', 'T')  # Draw the underline

    # Add Level (italic)
    pdf.ln(5)  # Space between title and level
    pdf.set_font('Arial', 'I', 12)
    pdf.cell(200, 10, f'Level: {level}', ln=True, align='C')

    # Add content from Groq response
    pdf.ln(10)  # Add a line break
    pdf.set_font('Arial', '', 12)
    
    # Split the response into lines and add to the PDF
    for line in cleaned_response.splitlines():
        pdf.multi_cell(0, 10, line)
    
    # Save the PDF to a file
    filename = f'roadmap.pdf'
    pdf.output(filename)

    print(f"PDF generated and saved as {filename}")
    return filename



@app.route('/mcq_test/<topic>/<subtopic>', methods=['GET'])
@login_required
def mcq_test(topic, subtopic):
    # Generate your questions as before
    questions = generate_mcq_questions(topic, subtopic)
    
    # Build a map of string indices to the correct answer (e.g. { "0": "A)", "1": "C)", ... })
    correct_answers = { str(i): q['answer'] for i, q in enumerate(questions) }

    return render_template(
        'mcq_test.html',
        topic=topic,
        subtopic=subtopic,
        questions=questions,
        correct_answers=json.dumps(correct_answers) 
    )

def generate_mcq_questions(topic, subtopic):
    prompt = f"Create 4 multiple-choice questions about {subtopic} under the topic {topic}. Include 4 options for each question and specify the correct answer."
    
    chat_completion = groq.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama3-8b-8192",
    )
    
    response = chat_completion.choices[0].message.content
    print(f"LLM Response: {response}")  # Debug print

    questions = parse_mcq_response(response)
    return questions

def parse_mcq_response(response):
    questions = []
    
    items = response.strip().splitlines()  # Split by lines

    current_question = None
    options = []
    correct_answer = None

    for index,line in enumerate(items):
        line = line.strip()
        temp="*Correct answer:*"
        temp2="Correct answer:"
        # Check for the start of a question
        line.strip("*")
        #print(type(line))
        line=line.replace("*","")
        print(line)
        if line.startswith("Question" or " Question"):
            if current_question is not None:  # Save the previous question before starting a new one
                questions.append({
                    'question': current_question,
                    'options': options,
                    'answer': correct_answer
                })

            current_question=line
            if(line=="Question 1" or line=="Question 2" or line=="Question 3" or line==" Question 1" or line==" Question 2" or line==" Question 3" or line==" Question 1:" or line==" Question 2:" or line==" Question 3:" or line=="Question 1:" or line=="Question 2:" or line=="Question 3:" or line=="Question 4:"  or line==" Question 4:"  or line=="Question 4: "  or line=="Question 4 :"):
                current_question=items[index+1]
            options = []  # Reset options
            correct_answer = None  # Reset correct answer
        
        # Check for options
        elif line.startswith(('A)', 'B)', 'C)', 'D)','a)','b)','c)','d)')):
            options.append(line)
        
        # Check for the correct answer
        elif line.startswith("Correct"):
            correct_answer = line # Extract the correct answer
    # Append the last question after finishing the loop
    if current_question is not None:
        questions.append({
            'question': current_question,
            'options': options,
            'answer': correct_answer
        })
    print(questions)
    return questions


if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)